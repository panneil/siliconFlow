"""
文档处理模块
支持文本、图像、Office、PDF等多种格式
"""

import os
import io
import re
import json
import base64
from PIL import Image
import pytesseract
import markdown
import PyPDF2
try:
    import docx
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentProcessor:
    """文档处理类"""
    
    def __init__(self, tesseract_path=None):
        """初始化文档处理器"""
        self.tesseract_path = tesseract_path
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
            
    def process_file(self, file_path):
        """处理文件"""
        # 获取文件扩展名
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # 根据文件类型调用不同的处理方法
        if ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.csv']:
            return self.process_text_file(file_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
            return self.process_image_file(file_path)
        elif ext == '.pdf':
            return self.process_pdf_file(file_path)
        elif ext == '.docx':
            return self.process_docx_file(file_path)
        elif ext == '.xlsx':
            return self.process_excel_file(file_path)
        elif ext == '.pptx':
            return self.process_pptx_file(file_path)
        else:
            return {"error": f"不支持的文件类型: {ext}"}
            
    def process_text_file(self, file_path):
        """处理文本文件"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                
            # 获取文件扩展名
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
            
            # 如果是Markdown文件，转换为HTML
            if ext == '.md':
                html_content = markdown.markdown(content)
                return {
                    "type": "markdown",
                    "content": content,
                    "html": html_content
                }
                
            # 如果是代码文件，添加语言标记
            if ext in ['.py', '.js', '.html', '.css', '.json']:
                language_map = {
                    '.py': 'python',
                    '.js': 'javascript',
                    '.html': 'html',
                    '.css': 'css',
                    '.json': 'json'
                }
                return {
                    "type": "code",
                    "language": language_map.get(ext, 'text'),
                    "content": content
                }
                
            # 默认作为普通文本处理
            return {
                "type": "text",
                "content": content
            }
        except Exception as e:
            return {"error": str(e)}
            
    def process_image_file(self, file_path):
        """处理图像文件"""
        try:
            # 读取图像
            img = Image.open(file_path)
            
            # 获取图像信息
            width, height = img.size
            format_name = img.format
            
            # 进行OCR识别
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            
            # 将图像转换为base64，方便显示
            buffered = io.BytesIO()
            img.save(buffered, format=format_name)
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            
            return {
                "type": "image",
                "width": width,
                "height": height,
                "format": format_name,
                "text": text,
                "base64": f"data:image/{format_name.lower()};base64,{img_base64}"
            }
        except Exception as e:
            return {"error": str(e)}
            
    def process_pdf_file(self, file_path):
        """处理PDF文件"""
        try:
            text_content = ""
            
            # 打开PDF文件
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                # 获取页数
                num_pages = len(pdf_reader.pages)
                
                # 提取文本
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text_content += page.extract_text() + "\n\n"
                    
            return {
                "type": "pdf",
                "num_pages": num_pages,
                "content": text_content
            }
        except Exception as e:
            return {"error": str(e)}
            
    def process_docx_file(self, file_path):
        """处理Word文档"""
        if not DOCX_AVAILABLE:
            return {"error": "python-docx未安装，无法处理Word文档"}
            
        try:
            doc = docx.Document(file_path)
            
            # 提取文本
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
                
            return {
                "type": "docx",
                "content": "\n".join(full_text)
            }
        except Exception as e:
            return {"error": str(e)}
            
    def process_excel_file(self, file_path):
        """处理Excel文件"""
        if not EXCEL_AVAILABLE:
            return {"error": "openpyxl未安装，无法处理Excel文件"}
            
        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
            
            result = {
                "type": "excel",
                "sheets": {}
            }
            
            # 处理每个工作表
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # 转换为二维数组
                data = []
                for row in sheet.iter_rows(values_only=True):
                    data.append(list(row))
                    
                result["sheets"][sheet_name] = data
                
            return result
        except Exception as e:
            return {"error": str(e)}
            
    def process_pptx_file(self, file_path):
        """处理PowerPoint文件"""
        if not PPTX_AVAILABLE:
            return {"error": "python-pptx未安装，无法处理PowerPoint文件"}
            
        try:
            prs = Presentation(file_path)
            
            result = {
                "type": "pptx",
                "slides": []
            }
            
            # 处理每个幻灯片
            for slide in prs.slides:
                slide_text = []
                
                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                        
                result["slides"].append("\n".join(slide_text))
                
            return result
        except Exception as e:
            return {"error": str(e)}
            
    def extract_text_from_image(self, image_path):
        """从图像中提取文本"""
        try:
            img = Image.open(image_path)
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            return text
        except Exception as e:
            return f"提取文本失败: {str(e)}"
            
    def convert_markdown_to_html(self, markdown_text):
        """将Markdown转换为HTML"""
        try:
            html = markdown.markdown(markdown_text)
            return html
        except Exception as e:
            return f"转换失败: {str(e)}"
            
    def render_mermaid_chart(self, mermaid_code):
        """渲染Mermaid图表"""
        # 返回Mermaid代码，让前端渲染
        return {
            "type": "mermaid",
            "code": mermaid_code
        }
        
    def process_clipboard_image(self, image_data):
        """处理剪贴板图像"""
        try:
            # 将图像数据转换为PIL Image对象
            image = Image.open(io.BytesIO(image_data))
            
            # 进行OCR识别
            text = pytesseract.image_to_string(image, lang='chi_sim+eng')
            
            # 将图像转换为base64，方便显示
            buffered = io.BytesIO()
            image.save(buffered, format="PNG")
            img_base64 = base64.b64encode(buffered.getvalue()).decode()
            
            return {
                "type": "clipboard_image",
                "width": image.width,
                "height": image.height,
                "text": text,
                "base64": f"data:image/png;base64,{img_base64}"
            }
        except Exception as e:
            return {"error": str(e)} 